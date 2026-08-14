"""
AI PPT Generator - Backend

Flask backend for generating professional PowerPoint presentations.

Flow:

User Topic
    ↓
Gemini Content Strategist
    ↓
Presentation Storyline
    ↓
Slide Type Selection
    ↓
Content + Visual Structure
    ↓
PowerPoint Layout Engine
    ↓
Editable PPTX
"""

import os
import io
import json
import random
import urllib.parse
import requests

from flask import Flask, request, send_file, jsonify
from flask_cors import CORS

from google import genai
from google.genai import types

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE, MSO_CONNECTOR

from PIL import Image, ImageDraw, ImageFilter


# ============================================================
# APP
# ============================================================

app = Flask(__name__)
CORS(app)


# ============================================================
# GEMINI
# ============================================================

# You can override this from CMD:
#
# set GEMINI_MODEL=your-model-name
#
# Otherwise the current default is used.

MODEL = os.environ.get(
    "GEMINI_MODEL",
    "gemini-3.6-flash"
)

_client = None


def get_client():

    global _client

    if _client is None:

        api_key = os.environ.get(
            "GEMINI_API_KEY"
        )

        if not api_key:

            raise RuntimeError(
                "GEMINI_API_KEY is not set. "
                "Set it before generating a presentation."
            )

        _client = genai.Client(
            api_key=api_key
        )

    return _client


# ============================================================
# THEMES
# ============================================================

THEMES = {

    "midnight": {
        "name": "Midnight",
        "bg": "0F1220",
        "primary": "6C5CE7",
        "secondary": "00CEC9",
        "text": "F5F6FA",
        "muted": "A0A3BD",
        "font": "Poppins",
        "gradient": [
            "1A1C2E",
            "2D2F4A"
        ],
    },

    "sunrise": {
        "name": "Sunrise",
        "bg": "FFF8F0",
        "primary": "FF6B35",
        "secondary": "F7C548",
        "text": "2B2118",
        "muted": "8A7968",
        "font": "Poppins",
        "gradient": [
            "FFDDB3",
            "FFB570"
        ],
    },

    "forest": {
        "name": "Forest",
        "bg": "F4F9F4",
        "primary": "1B6B3E",
        "secondary": "8FBF6E",
        "text": "18291D",
        "muted": "5F7A65",
        "font": "Poppins",
        "gradient": [
            "C9E4CA",
            "8FBF6E"
        ],
    },

    "ocean": {
        "name": "Ocean",
        "bg": "F0F7FF",
        "primary": "0B6E99",
        "secondary": "37B6E8",
        "text": "0B1F2A",
        "muted": "5C7A8A",
        "font": "Poppins",
        "gradient": [
            "BFE3F5",
            "6FC3E8"
        ],
    },

    "mono": {
        "name": "Monochrome",
        "bg": "FFFFFF",
        "primary": "1A1A1A",
        "secondary": "6E6E6E",
        "text": "111111",
        "muted": "8C8C8C",
        "font": "Poppins",
        "gradient": [
            "E8E8E8",
            "CFCFCF"
        ],
    },
}


# ============================================================
# BASIC HELPERS
# ============================================================

def hexc(value):

    return RGBColor.from_string(
        value
    )


def safe_text(value):

    if value is None:
        return ""

    return str(value).strip()


# ============================================================
# IMAGE GENERATION
# ============================================================

def image_url_for(
    query,
    theme,
    w=1000,
    h=1000
):

    prompt = (
        f"{query}, "
        f"{theme['name']} color palette, "
        f"professional presentation visual, "
        f"modern corporate design, "
        f"high quality, "
        f"no text, "
        f"no watermark"
    )

    encoded = urllib.parse.quote(
        prompt
    )

    seed = abs(
        hash(query)
    ) % 100000

    return (
        "https://image.pollinations.ai/prompt/"
        f"{encoded}"
        f"?width={w}"
        f"&height={h}"
        f"&seed={seed}"
        f"&nologo=true"
    )


def fetch_image_bytes(
    url,
    timeout=25
):

    response = requests.get(
        url,
        timeout=timeout
    )

    response.raise_for_status()

    return io.BytesIO(
        response.content
    )


def generate_art(
    theme,
    seed_text,
    w=1200,
    h=900
):

    c1 = tuple(
        int(
            theme["gradient"][0][i:i + 2],
            16
        )
        for i in (0, 2, 4)
    )

    c2 = tuple(
        int(
            theme["gradient"][1][i:i + 2],
            16
        )
        for i in (0, 2, 4)
    )

    img = Image.new(
        "RGB",
        (w, h),
        c1
    )

    draw = ImageDraw.Draw(
        img
    )

    for y in range(h):

        t = y / h

        r = int(
            c1[0] +
            (c2[0] - c1[0]) * t
        )

        g = int(
            c1[1] +
            (c2[1] - c1[1]) * t
        )

        b = int(
            c1[2] +
            (c2[2] - c1[2]) * t
        )

        draw.line(
            [(0, y), (w, y)],
            fill=(r, g, b)
        )

    rnd = random.Random(
        seed_text
    )

    accent = tuple(
        int(
            theme["primary"][i:i + 2],
            16
        )
        for i in (0, 2, 4)
    )

    for _ in range(6):

        x = rnd.randint(
            0,
            w
        )

        y = rnd.randint(
            0,
            h
        )

        radius = rnd.randint(
            60,
            220
        )

        overlay = Image.new(
            "RGBA",
            (w, h),
            (0, 0, 0, 0)
        )

        od = ImageDraw.Draw(
            overlay
        )

        od.ellipse(
            [
                x - radius,
                y - radius,
                x + radius,
                y + radius
            ],
            fill=accent + (35,)
        )

        img = Image.alpha_composite(
            img.convert("RGBA"),
            overlay
        ).convert("RGB")

    img = img.filter(
        ImageFilter.GaussianBlur(2)
    )

    buffer = io.BytesIO()

    img.save(
        buffer,
        format="PNG"
    )

    buffer.seek(0)

    return buffer


# ============================================================
# AI PRESENTATION STRATEGY PROMPT
# ============================================================

OUTLINE_SYSTEM_PROMPT = """
You are an expert presentation strategist, researcher,
storytelling designer and visual communication expert.

Your job is NOT to simply generate bullet points.

You must create a professional presentation with:

1. Strong logical storyline
2. Progressive flow between slides
3. Meaningful content
4. Clear key message on every slide
5. Different visual layouts
6. Appropriate diagrams
7. Concise professional language
8. No repetitive slides

Return ONLY valid JSON.

DO NOT use markdown.
DO NOT use ```json.
DO NOT add explanations outside JSON.


============================================================
JSON STRUCTURE
============================================================

{
  "title": "Presentation title",
  "subtitle": "Short subtitle",
  "storyline": "One sentence describing the overall story",

  "slides": [

    {
      "type": "title",
      "title": "",
      "purpose": "",
      "key_message": "",
      "bullets": [],
      "visual_type": "none",
      "needs_image": false,
      "image_query": "",
      "flowchart_steps": [],
      "causes": [],
      "effects": [],
      "comparison": {
        "left_title": "",
        "right_title": "",
        "left_points": [],
        "right_points": []
      },
      "timeline": []
    }

  ]
}


============================================================
ALLOWED SLIDE TYPES
============================================================

title

key_stat

content

comparison

process

cause_effect

hierarchy

cycle

timeline

image

solutions

closing


============================================================
STORYLINE RULES
============================================================

Every presentation must tell a story.

Do NOT make slides independent from each other.

For general topics prefer:

Introduction
→ Context
→ Problem
→ Causes
→ How it works
→ Impact
→ Solutions
→ Conclusion


For technical topics prefer:

What is it?
→ Why does it matter?
→ How does it work?
→ Architecture / Process
→ Example
→ Advantages
→ Limitations
→ Future


For business topics prefer:

Market
→ Problem
→ Customer
→ Solution
→ Business Model
→ Competition
→ Strategy
→ Conclusion


For academic topics prefer:

Introduction
→ Background
→ Core Concept
→ Methodology
→ Analysis
→ Results
→ Applications
→ Conclusion


============================================================
CONTENT RULES
============================================================

Every slide must have ONE clear key message.

Avoid generic statements.

BAD:

"Technology is changing the world."

GOOD:

"Automation reduces repetitive work while increasing
demand for digital skills."

Avoid repetition.

Do not repeat the definition of the topic on multiple slides.

Use examples whenever useful.

Prefer concrete facts, relationships,
examples and meaningful explanations.

Normal content slides should have maximum 3-4 bullets.

Bullets should generally be short.

Do NOT make every slide a bullet slide.


============================================================
VISUAL RULES
============================================================

Use visual structures when they improve understanding.

Use:

process
for sequential steps.

cause_effect
for causes leading to a central problem
and resulting effects.

hierarchy
for parent-child relationships.

cycle
for feedback loops and repeating systems.

comparison
for comparing two concepts.

timeline
for chronological development.

key_stat
for important numbers or statistics.

image
when a real-world visual adds value.

solutions
for recommendations or action plans.


============================================================
PROCESS RULES
============================================================

For a process:

Input
→ Processing
→ Decision
→ Output

OR

Step A
→ Step B
→ Step C
→ Step D


============================================================
CAUSE/EFFECT RULES
============================================================

For cause_effect:

causes → central problem → effects

Example:

Economic slowdown
        ↓
Skills mismatch → Unemployment ← Automation
        ↓
Lower income
        ↓
Social impact


============================================================
HIERARCHY RULES
============================================================

Use:

Main concept
→ Category
→ Subcategory


============================================================
CYCLE RULES
============================================================

The final stage should connect back
to the first stage.

A → B → C → D → A


============================================================
COMPARISON RULES
============================================================

For comparison slides provide:

left_title

right_title

left_points

right_points

Use when two concepts can meaningfully
be compared.


============================================================
TIMELINE RULES
============================================================

Use chronological events.

Each timeline item should contain:

year/date
event
short explanation


============================================================
FLOWCHART QUALITY
============================================================

Do NOT create a flowchart simply to fill space.

Do NOT create generic:

Step 1 → Step 2 → Step 3

unless the subject is actually sequential.

Prefer meaningful relationships.

Use diagrams to explain concepts,
not decoration.


============================================================
SLIDE COUNT
============================================================

Generate approximately the requested number
of slides.

Always include:

1 title slide

1 conclusion slide

Body slides should form a logical narrative.

Do not create unnecessary slides.


============================================================
FINAL QUALITY CHECK
============================================================

Before returning JSON verify:

- Storyline is logical
- No repeated slides
- No generic filler
- No unnecessary flowcharts
- Visual types are appropriate
- Content is concise
- Each slide has one key message
- Title and conclusion exist
- JSON is valid

Return ONLY JSON.
"""


# ============================================================
# GENERATE OUTLINE
# ============================================================

def generate_outline(
    topic,
    num_slides,
    audience
):

    user_prompt = f"""
Topic:
{topic}

Target slide count:
{num_slides}

Audience:
{audience or "general audience"}

Create a professional presentation.

The presentation should feel like it was designed
by an experienced presentation consultant.

Do not merely summarize the topic.

Build a clear story and choose the appropriate
visual structure for every slide.

Generate the final JSON now.
"""

    response = get_client().models.generate_content(

        model=MODEL,

        contents=user_prompt,

        config=types.GenerateContentConfig(

            system_instruction=
            OUTLINE_SYSTEM_PROMPT,

            response_mime_type=
            "application/json",

            max_output_tokens=5000
        )
    )

    raw = response.text.strip()

    raw = (
        raw
        .replace(
            "```json",
            ""
        )
        .replace(
            "```",
            ""
        )
        .strip()
    )

    outline = json.loads(
        raw
    )

    return outline


# ============================================================
# IMAGE URL ATTACHMENT
# ============================================================

def attach_image_urls(
    outline,
    theme_key
):

    theme = THEMES.get(
        theme_key,
        THEMES["midnight"]
    )

    for slide in outline.get(
        "slides",
        []
    ):

        if (
            slide.get("needs_image")
            and slide.get("image_query")
        ):

            slide["image_url"] = (
                image_url_for(
                    slide["image_query"],
                    theme
                )
            )

    return outline


# ============================================================
# PPT HELPERS
# ============================================================

def add_text(
    slide,
    text,
    left,
    top,
    width,
    height,
    size,
    color,
    theme,
    bold=False,
    align=PP_ALIGN.LEFT
):

    box = slide.shapes.add_textbox(
        left,
        top,
        width,
        height
    )

    tf = box.text_frame

    tf.clear()

    tf.word_wrap = True

    tf.vertical_anchor = (
        MSO_ANCHOR.MIDDLE
    )

    p = tf.paragraphs[0]

    p.alignment = align

    run = p.add_run()

    run.text = safe_text(text)

    run.font.size = Pt(size)

    run.font.bold = bold

    run.font.name = theme["font"]

    run.font.color.rgb = (
        hexc(color)
    )

    return box


def add_header(
    slide,
    title,
    theme
):

    add_text(
        slide,
        title,
        Inches(0.7),
        Inches(0.45),
        Inches(11.9),
        Inches(0.65),
        27,
        theme["primary"],
        theme,
        bold=True
    )

    rule = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(0.7),
        Inches(1.22),
        Inches(1.8),
        Pt(3)
    )

    rule.fill.solid()

    rule.fill.fore_color.rgb = (
        hexc(theme["secondary"])
    )

    rule.line.fill.background()


def add_footer(
    slide,
    slide_number,
    theme
):

    add_text(
        slide,
        str(slide_number),
        Inches(12.2),
        Inches(7.05),
        Inches(0.5),
        Inches(0.25),
        9,
        theme["muted"],
        theme,
        align=PP_ALIGN.RIGHT
    )


def add_card(
    slide,
    title,
    body,
    left,
    top,
    width,
    height,
    theme,
    accent=None
):

    if accent is None:
        accent = theme["primary"]

    shape = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        left,
        top,
        width,
        height
    )

    shape.fill.solid()

    shape.fill.fore_color.rgb = (
        hexc(theme["bg"])
    )

    shape.line.color.rgb = (
        hexc(theme["muted"])
    )

    shape.line.width = Pt(1)

    add_text(
        slide,
        title,
        left + Inches(0.18),
        top + Inches(0.15),
        width - Inches(0.36),
        Inches(0.4),
        15,
        accent,
        theme,
        bold=True
    )

    add_text(
        slide,
        body,
        left + Inches(0.18),
        top + Inches(0.62),
        width - Inches(0.36),
        height - Inches(0.78),
        12,
        theme["text"],
        theme
    )

    return shape


# ============================================================
# PROCESS DIAGRAM
# ============================================================

def add_process_diagram(
    slide,
    steps,
    theme,
    left,
    top,
    width,
    height
):

    steps = [
        safe_text(x)
        for x in steps
        if safe_text(x)
    ]

    if not steps:
        return

    count = min(
        len(steps),
        5
    )

    steps = steps[:count]

    gap = Inches(0.28)

    box_width = (
        width -
        gap * (count - 1)
    ) / count

    box_height = Inches(1.15)

    y = top + (
        height - box_height
    ) / 2

    boxes = []

    for i, step in enumerate(steps):

        x = (
            left +
            i * (box_width + gap)
        )

        box = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE,
            x,
            y,
            box_width,
            box_height
        )

        box.fill.solid()

        box.fill.fore_color.rgb = (
            hexc(
                theme["primary"]
                if i % 2 == 0
                else theme["secondary"]
            )
        )

        box.line.fill.background()

        tf = box.text_frame

        tf.clear()

        tf.word_wrap = True

        tf.vertical_anchor = (
            MSO_ANCHOR.MIDDLE
        )

        p = tf.paragraphs[0]

        p.alignment = PP_ALIGN.CENTER

        run = p.add_run()

        run.text = step

        run.font.size = Pt(13)

        run.font.bold = True

        run.font.name = theme["font"]

        run.font.color.rgb = (
            hexc(
                theme["bg"]
                if i % 2 == 0
                else theme["text"]
            )
        )

        boxes.append(box)

    for i in range(
        len(boxes) - 1
    ):

        connector = (
            slide.shapes.add_connector(
                MSO_CONNECTOR.STRAIGHT,

                boxes[i].left +
                boxes[i].width,

                boxes[i].top +
                boxes[i].height // 2,

                boxes[i + 1].left,

                boxes[i + 1].top +
                boxes[i + 1].height // 2
            )
        )

        connector.line.color.rgb = (
            hexc(theme["muted"])
        )

        connector.line.width = Pt(2)


# ============================================================
# CAUSE / EFFECT DIAGRAM
# ============================================================

def add_cause_effect(
    slide,
    causes,
    effects,
    central_problem,
    theme
):

    causes = [
        safe_text(x)
        for x in causes
        if safe_text(x)
    ][:4]

    effects = [
        safe_text(x)
        for x in effects
        if safe_text(x)
    ][:4]

    center_x = Inches(5.0)
    center_y = Inches(3.0)

    center_w = Inches(3.3)
    center_h = Inches(1.2)

    center = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        center_x,
        center_y,
        center_w,
        center_h
    )

    center.fill.solid()

    center.fill.fore_color.rgb = (
        hexc(theme["primary"])
    )

    center.line.fill.background()

    tf = center.text_frame

    tf.clear()

    tf.vertical_anchor = (
        MSO_ANCHOR.MIDDLE
    )

    p = tf.paragraphs[0]

    p.alignment = PP_ALIGN.CENTER

    run = p.add_run()

    run.text = (
        central_problem
        or "Central Problem"
    )

    run.font.size = Pt(17)

    run.font.bold = True

    run.font.color.rgb = (
        hexc(theme["bg"])
    )

    # Causes

    for i, cause in enumerate(causes):

        y = (
            Inches(1.65)
            + i * Inches(1.05)
        )

        card = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE,
            Inches(0.7),
            y,
            Inches(3.4),
            Inches(0.72)
        )

        card.fill.solid()

        card.fill.fore_color.rgb = (
            hexc(theme["secondary"])
        )

        card.line.fill.background()

        tf = card.text_frame

        tf.clear()

        tf.vertical_anchor = (
            MSO_ANCHOR.MIDDLE
        )

        p = tf.paragraphs[0]

        p.alignment = PP_ALIGN.CENTER

        r = p.add_run()

        r.text = cause

        r.font.size = Pt(12)

        r.font.bold = True

        r.font.color.rgb = (
            hexc(theme["text"])
        )

        connector = (
            slide.shapes.add_connector(
                MSO_CONNECTOR.STRAIGHT,
                Inches(4.1),
                y + Inches(0.36),
                center_x,
                center_y +
                center_h / 2
            )
        )

        connector.line.color.rgb = (
            hexc(theme["muted"])
        )

    # Effects

    for i, effect in enumerate(effects):

        y = (
            Inches(1.65)
            + i * Inches(1.05)
        )

        card = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE,
            Inches(9.2),
            y,
            Inches(3.4),
            Inches(0.72)
        )

        card.fill.solid()

        card.fill.fore_color.rgb = (
            hexc(theme["primary"])
        )

        card.line.fill.background()

        tf = card.text_frame

        tf.clear()

        tf.vertical_anchor = (
            MSO_ANCHOR.MIDDLE
        )

        p = tf.paragraphs[0]

        p.alignment = PP_ALIGN.CENTER

        r = p.add_run()

        r.text = effect

        r.font.size = Pt(12)

        r.font.bold = True

        r.font.color.rgb = (
            hexc(theme["bg"])
        )

        connector = (
            slide.shapes.add_connector(
                MSO_CONNECTOR.STRAIGHT,
                center_x + center_w,
                center_y +
                center_h / 2,
                Inches(9.2),
                y + Inches(0.36)
            )
        )

        connector.line.color.rgb = (
            hexc(theme["muted"])
        )


# ============================================================
# HIERARCHY DIAGRAM
# ============================================================

def add_hierarchy(
    slide,
    title,
    categories,
    theme
):

    root = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        Inches(4.7),
        Inches(1.8),
        Inches(3.9),
        Inches(0.85)
    )

    root.fill.solid()

    root.fill.fore_color.rgb = (
        hexc(theme["primary"])
    )

    root.line.fill.background()

    tf = root.text_frame

    tf.clear()

    tf.vertical_anchor = (
        MSO_ANCHOR.MIDDLE
    )

    p = tf.paragraphs[0]

    p.alignment = PP_ALIGN.CENTER

    r = p.add_run()

    r.text = title

    r.font.size = Pt(16)

    r.font.bold = True

    r.font.color.rgb = (
        hexc(theme["bg"])
    )

    categories = categories[:4]

    if not categories:
        return

    gap = Inches(0.25)

    total_width = Inches(11.5)

    card_width = (
        total_width -
        gap * (len(categories) - 1)
    ) / len(categories)

    start_x = Inches(0.9)

    for i, category in enumerate(
        categories
    ):

        if isinstance(
            category,
            dict
        ):

            name = category.get(
                "name",
                ""
            )

            children = category.get(
                "children",
                []
            )

        else:

            name = str(category)

            children = []

        x = (
            start_x +
            i * (
                card_width + gap
            )
        )

        card = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE,
            x,
            Inches(3.45),
            card_width,
            Inches(1.0)
        )

        card.fill.solid()

        card.fill.fore_color.rgb = (
            hexc(theme["secondary"])
        )

        card.line.fill.background()

        tf = card.text_frame

        tf.clear()

        tf.vertical_anchor = (
            MSO_ANCHOR.MIDDLE
        )

        p = tf.paragraphs[0]

        p.alignment = PP_ALIGN.CENTER

        r = p.add_run()

        r.text = name

        r.font.size = Pt(14)

        r.font.bold = True

        r.font.color.rgb = (
            hexc(theme["text"])
        )

        connector = (
            slide.shapes.add_connector(
                MSO_CONNECTOR.STRAIGHT,
                Inches(6.65),
                Inches(2.65),
                x + card_width / 2,
                Inches(3.45)
            )
        )

        connector.line.color.rgb = (
            hexc(theme["muted"])
        )

        if children:

            child_text = "\n".join(
                [
                    f"• {safe_text(child)}"
                    for child in children[:3]
                ]
            )

            add_text(
                slide,
                child_text,
                x,
                Inches(4.65),
                card_width,
                Inches(1.2),
                11,
                theme["text"],
                theme
            )


# ============================================================
# CYCLE DIAGRAM
# ============================================================

def add_cycle(
    slide,
    steps,
    theme
):

    steps = [
        safe_text(x)
        for x in steps
        if safe_text(x)
    ][:5]

    if len(steps) < 2:
        return

    positions = [

        (
            Inches(5.1),
            Inches(1.65)
        ),

        (
            Inches(8.5),
            Inches(3.0)
        ),

        (
            Inches(5.1),
            Inches(4.65)
        ),

        (
            Inches(1.7),
            Inches(3.0)
        ),

        (
            Inches(5.1),
            Inches(3.0)
        )
    ]

    boxes = []

    for i, step in enumerate(
        steps
    ):

        x, y = positions[i]

        width = Inches(2.8)

        height = Inches(0.9)

        box = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE,
            x,
            y,
            width,
            height
        )

        box.fill.solid()

        box.fill.fore_color.rgb = (
            hexc(
                theme["primary"]
                if i % 2 == 0
                else theme["secondary"]
            )
        )

        box.line.fill.background()

        tf = box.text_frame

        tf.clear()

        tf.vertical_anchor = (
            MSO_ANCHOR.MIDDLE
        )

        p = tf.paragraphs[0]

        p.alignment = PP_ALIGN.CENTER

        r = p.add_run()

        r.text = step

        r.font.size = Pt(12)

        r.font.bold = True

        r.font.color.rgb = (
            hexc(theme["bg"])
        )

        boxes.append(box)

    for i in range(
        len(boxes)
    ):

        current = boxes[i]

        nxt = boxes[
            (i + 1) %
            len(boxes)
        ]

        x1 = (
            current.left +
            current.width / 2
        )

        y1 = (
            current.top +
            current.height / 2
        )

        x2 = (
            nxt.left +
            nxt.width / 2
        )

        y2 = (
            nxt.top +
            nxt.height / 2
        )

        connector = (
            slide.shapes.add_connector(
                MSO_CONNECTOR.STRAIGHT,
                x1,
                y1,
                x2,
                y2
            )
        )

        connector.line.color.rgb = (
            hexc(theme["muted"])
        )

        connector.line.width = Pt(2)


# ============================================================
# COMPARISON
# ============================================================

def add_comparison(
    slide,
    comparison,
    theme
):

    left_title = comparison.get(
        "left_title",
        "Option A"
    )

    right_title = comparison.get(
        "right_title",
        "Option B"
    )

    left_points = comparison.get(
        "left_points",
        []
    )

    right_points = comparison.get(
        "right_points",
        []
    )

    add_card(
        slide,
        left_title,
        "\n".join(
            [
                f"• {safe_text(x)}"
                for x in left_points[:5]
            ]
        ),
        Inches(0.8),
        Inches(2.0),
        Inches(5.7),
        Inches(4.2),
        theme,
        theme["primary"]
    )

    add_card(
        slide,
        right_title,
        "\n".join(
            [
                f"• {safe_text(x)}"
                for x in right_points[:5]
            ]
        ),
        Inches(6.85),
        Inches(2.0),
        Inches(5.7),
        Inches(4.2),
        theme,
        theme["secondary"]
    )


# ============================================================
# KEY STAT
# ============================================================

def add_key_stat(
    slide,
    slide_data,
    theme
):

    key_message = safe_text(
        slide_data.get(
            "key_message",
            ""
        )
    )

    bullets = slide_data.get(
        "bullets",
        []
    )

    statistic = ""

    if bullets:

        statistic = safe_text(
            bullets[0]
        )

    add_text(
        slide,
        statistic,
        Inches(0.8),
        Inches(2.0),
        Inches(11.7),
        Inches(1.5),
        42,
        theme["primary"],
        theme,
        bold=True,
        align=PP_ALIGN.CENTER
    )

    add_text(
        slide,
        key_message,
        Inches(1.5),
        Inches(3.8),
        Inches(10.3),
        Inches(1.4),
        20,
        theme["text"],
        theme,
        align=PP_ALIGN.CENTER
    )


# ============================================================
# TIMELINE
# ============================================================

def add_timeline(
    slide,
    timeline,
    theme
):

    timeline = timeline[:5]

    if not timeline:
        return

    line_y = Inches(3.5)

    connector = (
        slide.shapes.add_connector(
            MSO_CONNECTOR.STRAIGHT,
            Inches(1.0),
            line_y,
            Inches(12.2),
            line_y
        )
    )

    connector.line.color.rgb = (
        hexc(theme["muted"])
    )

    connector.line.width = Pt(3)

    spacing = (
        Inches(11.0) /
        max(
            len(timeline) - 1,
            1
        )
    )

    for i, item in enumerate(
        timeline
    ):

        if isinstance(
            item,
            dict
        ):

            date = safe_text(
                item.get(
                    "date",
                    item.get(
                        "year",
                        ""
                    )
                )
            )

            event = safe_text(
                item.get(
                    "event",
                    ""
                )
            )

            description = safe_text(
                item.get(
                    "description",
                    ""
                )
            )

        else:

            date = ""

            event = safe_text(
                item
            )

            description = ""

        x = Inches(1.0) + (
            spacing * i
        )

        circle = slide.shapes.add_shape(
            MSO_SHAPE.OVAL,
            x - Inches(0.16),
            line_y - Inches(0.16),
            Inches(0.32),
            Inches(0.32)
        )

        circle.fill.solid()

        circle.fill.fore_color.rgb = (
            hexc(theme["primary"])
        )

        circle.line.fill.background()

        add_text(
            slide,
            date,
            x - Inches(0.7),
            Inches(2.55),
            Inches(1.4),
            Inches(0.4),
            12,
            theme["primary"],
            theme,
            bold=True,
            align=PP_ALIGN.CENTER
        )

        add_text(
            slide,
            event,
            x - Inches(1.0),
            Inches(3.9),
            Inches(2.0),
            Inches(0.6),
            12,
            theme["text"],
            theme,
            bold=True,
            align=PP_ALIGN.CENTER
        )

        if description:

            add_text(
                slide,
                description,
                x - Inches(1.0),
                Inches(4.55),
                Inches(2.0),
                Inches(1.0),
                10,
                theme["muted"],
                theme,
                align=PP_ALIGN.CENTER
            )


# ============================================================
# BULLET CONTENT
# ============================================================

def add_bullets(
    slide,
    bullets,
    theme,
    left=Inches(0.9),
    top=Inches(2.0),
    width=Inches(11.5),
    height=Inches(4.5)
):

    bullets = [
        safe_text(x)
        for x in bullets
        if safe_text(x)
    ][:4]

    box = slide.shapes.add_textbox(
        left,
        top,
        width,
        height
    )

    tf = box.text_frame

    tf.clear()

    tf.word_wrap = True

    for i, bullet in enumerate(
        bullets
    ):

        p = (
            tf.paragraphs[0]
            if i == 0
            else tf.add_paragraph()
        )

        p.space_after = Pt(18)

        r1 = p.add_run()

        r1.text = "●  "

        r1.font.size = Pt(16)

        r1.font.color.rgb = (
            hexc(theme["secondary"])
        )

        r2 = p.add_run()

        r2.text = bullet

        r2.font.size = Pt(16)

        r2.font.name = theme["font"]

        r2.font.color.rgb = (
            hexc(theme["text"])
        )


# ============================================================
# IMAGE SLIDE
# ============================================================

def add_image_slide(
    slide,
    slide_data,
    theme
):

    image_url = slide_data.get(
        "image_url"
    )

    try:

        if not image_url:
            raise ValueError(
                "Image URL missing"
            )

        image_buffer = (
            fetch_image_bytes(
                image_url
            )
        )

    except Exception:

        image_buffer = generate_art(
            theme,
            slide_data.get(
                "title",
                "visual"
            )
        )

    slide.shapes.add_picture(
        image_buffer,
        Inches(7.2),
        Inches(1.7),
        width=Inches(5.3),
        height=Inches(4.8)
    )

    add_text(
        slide,
        slide_data.get(
            "key_message",
            ""
        ),
        Inches(0.9),
        Inches(2.2),
        Inches(5.5),
        Inches(2.5),
        23,
        theme["text"],
        theme,
        bold=True
    )


# ============================================================
# SOLUTIONS SLIDE
# ============================================================

def add_solutions(
    slide,
    bullets,
    theme
):

    bullets = bullets[:4]

    if not bullets:
        return

    card_width = Inches(5.55)

    positions = [

        (
            Inches(0.8),
            Inches(2.0)
        ),

        (
            Inches(6.95),
            Inches(2.0)
        ),

        (
            Inches(0.8),
            Inches(4.25)
        ),

        (
            Inches(6.95),
            Inches(4.25)
        )
    ]

    for i, bullet in enumerate(
        bullets
    ):

        x, y = positions[i]

        add_card(
            slide,
            f"Action {i + 1}",
            bullet,
            x,
            y,
            card_width,
            Inches(1.75),
            theme,
            theme["primary"]
            if i % 2 == 0
            else theme["secondary"]
        )


# ============================================================
# BUILD PPT
# ============================================================

def build_pptx(
    outline,
    theme_key
):

    theme = THEMES.get(
        theme_key,
        THEMES["midnight"]
    )

    prs = Presentation()

    prs.slide_width = (
        Inches(13.333)
    )

    prs.slide_height = (
        Inches(7.5)
    )

    blank = prs.slide_layouts[6]


    # --------------------------------------------------------
    # NEW SLIDE
    # --------------------------------------------------------

    def new_slide():

        slide = prs.slides.add_slide(
            blank
        )

        bg = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            0,
            0,
            prs.slide_width,
            prs.slide_height
        )

        bg.fill.solid()

        bg.fill.fore_color.rgb = (
            hexc(theme["bg"])
        )

        bg.line.fill.background()

        bg.shadow.inherit = False

        slide.shapes._spTree.remove(
            bg._element
        )

        slide.shapes._spTree.insert(
            2,
            bg._element
        )

        return slide


    # --------------------------------------------------------
    # TITLE
    # --------------------------------------------------------

    slide = new_slide()

    accent = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        0,
        0,
        Inches(0.22),
        prs.slide_height
    )

    accent.fill.solid()

    accent.fill.fore_color.rgb = (
        hexc(theme["primary"])
    )

    accent.line.fill.background()

    add_text(
        slide,
        outline.get(
            "title",
            "Untitled Presentation"
        ),
        Inches(0.95),
        Inches(2.55),
        Inches(11.4),
        Inches(1.2),
        40,
        theme["text"],
        theme,
        bold=True
    )

    add_text(
        slide,
        outline.get(
            "subtitle",
            ""
        ),
        Inches(0.98),
        Inches(3.9),
        Inches(10.5),
        Inches(0.8),
        18,
        theme["muted"],
        theme
    )


    # --------------------------------------------------------
    # BODY SLIDES
    # --------------------------------------------------------

    body_slides = [
        s
        for s in outline.get(
            "slides",
            []
        )
        if s.get("type") not in (
            "title",
            "closing"
        )
    ]


    for slide_number, sl in enumerate(
        body_slides,
        start=2
    ):

        slide = new_slide()

        slide_type = sl.get(
            "type",
            "content"
        )

        title = safe_text(
            sl.get(
                "title",
                ""
            )
        )

        add_header(
            slide,
            title,
            theme
        )

        # ----------------------------------------------------
        # KEY STAT
        # ----------------------------------------------------

        if slide_type == "key_stat":

            add_key_stat(
                slide,
                sl,
                theme
            )


        # ----------------------------------------------------
        # PROCESS
        # ----------------------------------------------------

        elif (
            slide_type == "process"
            or sl.get(
                "visual_type"
            ) == "process"
        ):

            steps = sl.get(
                "flowchart_steps",
                []
            )

            add_process_diagram(
                slide,
                steps,
                theme,
                Inches(0.7),
                Inches(2.0),
                Inches(11.9),
                Inches(2.4)
            )

            add_text(
                slide,
                sl.get(
                    "key_message",
                    ""
                ),
                Inches(1.0),
                Inches(4.9),
                Inches(11.2),
                Inches(1.0),
                17,
                theme["text"],
                theme,
                align=PP_ALIGN.CENTER
            )


        # ----------------------------------------------------
        # CAUSE EFFECT
        # ----------------------------------------------------

        elif (
            slide_type == "cause_effect"
            or sl.get(
                "visual_type"
            ) == "cause_effect"
        ):

            add_cause_effect(
                slide,
                sl.get(
                    "causes",
                    []
                ),
                sl.get(
                    "effects",
                    []
                ),
                sl.get(
                    "key_message",
                    title
                ),
                theme
            )


        # ----------------------------------------------------
        # HIERARCHY
        # ----------------------------------------------------

        elif (
            slide_type == "hierarchy"
            or sl.get(
                "visual_type"
            ) == "hierarchy"
        ):

            categories = sl.get(
                "categories",
                []
            )

            add_hierarchy(
                slide,
                sl.get(
                    "key_message",
                    title
                ),
                categories,
                theme
            )


        # ----------------------------------------------------
        # CYCLE
        # ----------------------------------------------------

        elif (
            slide_type == "cycle"
            or sl.get(
                "visual_type"
            ) == "cycle"
        ):

            add_cycle(
                slide,
                sl.get(
                    "flowchart_steps",
                    []
                ),
                theme
            )


        # ----------------------------------------------------
        # COMPARISON
        # ----------------------------------------------------

        elif slide_type == "comparison":

            add_comparison(
                slide,
                sl.get(
                    "comparison",
                    {}
                ),
                theme
            )


        # ----------------------------------------------------
        # TIMELINE
        # ----------------------------------------------------

        elif slide_type == "timeline":

            add_timeline(
                slide,
                sl.get(
                    "timeline",
                    []
                ),
                theme
            )


        # ----------------------------------------------------
        # IMAGE
        # ----------------------------------------------------

        elif slide_type == "image":

            add_image_slide(
                slide,
                sl,
                theme
            )


        # ----------------------------------------------------
        # SOLUTIONS
        # ----------------------------------------------------

        elif slide_type == "solutions":

            add_solutions(
                slide,
                sl.get(
                    "bullets",
                    []
                ),
                theme
            )


        # ----------------------------------------------------
        # NORMAL CONTENT
        # ----------------------------------------------------

        else:

            image_required = (
                sl.get(
                    "needs_image"
                )
                and sl.get(
                    "image_url"
                )
            )

            if image_required:

                try:

                    image_buffer = (
                        fetch_image_bytes(
                            sl["image_url"]
                        )
                    )

                except Exception:

                    image_buffer = (
                        generate_art(
                            theme,
                            title
                        )
                    )

                slide.shapes.add_picture(
                    image_buffer,
                    Inches(8.2),
                    Inches(1.75),
                    width=Inches(4.3),
                    height=Inches(4.8)
                )

                add_text(
                    slide,
                    sl.get(
                        "key_message",
                        ""
                    ),
                    Inches(0.9),
                    Inches(1.8),
                    Inches(6.8),
                    Inches(1.2),
                    21,
                    theme["text"],
                    theme,
                    bold=True
                )

                add_bullets(
                    slide,
                    sl.get(
                        "bullets",
                        []
                    ),
                    theme,
                    Inches(0.9),
                    Inches(3.0),
                    Inches(6.8),
                    Inches(3.0)
                )

            else:

                key_message = safe_text(
                    sl.get(
                        "key_message",
                        ""
                    )
                )

                if key_message:

                    add_text(
                        slide,
                        key_message,
                        Inches(0.9),
                        Inches(1.75),
                        Inches(11.4),
                        Inches(1.1),
                        21,
                        theme["primary"],
                        theme,
                        bold=True
                    )

                    bullet_top = (
                        Inches(3.0)
                    )

                else:

                    bullet_top = (
                        Inches(2.0)
                    )

                add_bullets(
                    slide,
                    sl.get(
                        "bullets",
                        []
                    ),
                    theme,
                    Inches(0.9),
                    bullet_top,
                    Inches(11.4),
                    Inches(3.8)
                )


        add_footer(
            slide,
            slide_number,
            theme
        )


    # --------------------------------------------------------
    # CLOSING
    # --------------------------------------------------------

    slide = new_slide()

    add_text(
        slide,
        "Key Takeaways",
        Inches(1.0),
        Inches(1.0),
        Inches(11.3),
        Inches(0.8),
        30,
        theme["primary"],
        theme,
        bold=True,
        align=PP_ALIGN.CENTER
    )

    closing_bullets = []

    if body_slides:

        for sl in body_slides[-3:]:

            message = safe_text(
                sl.get(
                    "key_message",
                    ""
                )
            )

            if message:

                closing_bullets.append(
                    message
                )

    closing_bullets = closing_bullets[:3]

    if closing_bullets:

        add_bullets(
            slide,
            closing_bullets,
            theme,
            Inches(2.0),
            Inches(2.4),
            Inches(9.3),
            Inches(2.8)
        )

    add_text(
        slide,
        "Thank You",
        Inches(1.0),
        Inches(5.5),
        Inches(11.3),
        Inches(0.8),
        26,
        theme["text"],
        theme,
        bold=True,
        align=PP_ALIGN.CENTER
    )


    # --------------------------------------------------------
    # SAVE
    # --------------------------------------------------------

    buffer = io.BytesIO()

    prs.save(
        buffer
    )

    buffer.seek(0)

    return buffer


# ============================================================
# ROUTES
# ============================================================

@app.route(
    "/api/themes",
    methods=["GET"]
)
def get_themes():

    return jsonify({

        key: {

            "name": value["name"],

            "primary": value["primary"],

            "secondary": value["secondary"],

            "bg": value["bg"]
        }

        for key, value in THEMES.items()
    })


# ============================================================
# GENERATE OUTLINE
# ============================================================

@app.route(
    "/api/generate-outline",
    methods=["POST"]
)
def api_generate_outline():

    data = request.get_json(
        silent=True
    ) or {}

    topic = safe_text(
        data.get(
            "topic",
            ""
        )
    )

    theme_key = data.get(
        "theme",
        "midnight"
    )

    if not topic:

        return jsonify({
            "error": "Topic is required"
        }), 400

    try:

        outline = generate_outline(

            topic,

            data.get(
                "num_slides",
                8
            ),

            data.get(
                "audience",
                ""
            )
        )

        outline = attach_image_urls(
            outline,
            theme_key
        )

        return jsonify(
            outline
        )

    except Exception as e:

        print(
            "Outline generation error:",
            repr(e)
        )

        return jsonify({
            "error": str(e)
        }), 500


# ============================================================
# GENERATE PPTX
# ============================================================

@app.route(
    "/api/generate-pptx",
    methods=["POST"]
)
def api_generate_pptx():

    data = request.get_json(
        silent=True
    ) or {}

    outline = data.get(
        "outline"
    )

    theme_key = data.get(
        "theme",
        "midnight"
    )

    if not outline:

        return jsonify({
            "error": "Outline is required"
        }), 400

    try:

        buffer = build_pptx(
            outline,
            theme_key
        )

        title = safe_text(
            outline.get(
                "title",
                "presentation"
            )
        )

        filename = "".join(

            c

            for c in title

            if (
                c.isalnum()
                or c in " -_"
            )
        ).strip()

        if not filename:

            filename = "presentation"

        return send_file(

            buffer,

            as_attachment=True,

            download_name=(
                f"{filename}.pptx"
            ),

            mimetype=(
                "application/vnd.openxmlformats-officedocument."
                "presentationml.presentation"
            )
        )

    except Exception as e:

        print(
            "PPT generation error:",
            repr(e)
        )

        return jsonify({
            "error": str(e)
        }), 500


# ============================================================
# HEALTH CHECK
# ============================================================

@app.route(
    "/api/health",
    methods=["GET"]
)
def health():

    return jsonify({
        "status": "ok",
        "model": MODEL
    })


# ============================================================
# START
# ============================================================

if __name__ == "__main__":

    app.run(
        debug=True,
        port=5000
    )
